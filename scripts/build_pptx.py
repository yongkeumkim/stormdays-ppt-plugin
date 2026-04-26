"""
build_pptx.py — 슬라이드 계획 JSON + 스톰데이즈 템플릿으로 PPTX 생성

Usage: python build_pptx.py <slide_plan_json> <output_pptx> [--template <template_path>]

디자인 원칙: 템플릿 슬라이드를 복제해 콘텐츠만 교체. 1픽셀도 변경 없음.
"""

import sys
import json
import copy
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from lxml import etree
except ImportError:
    import subprocess
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "lxml", "-q"])
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from lxml import etree

# ── 네임스페이스 ──────────────────────────────────────────────────────────────
NS_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P  = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

def qn_a(tag): return f"{{{NS_A}}}{tag}"
def qn_p(tag): return f"{{{NS_P}}}{tag}"

# 템플릿 슬라이드 인덱스 (0-based)
TEMPLATE = {
    "cover":   0,   # slide1: 표지
    "toc":     1,   # slide2: 목차
    "stat":    2,   # slide3: 큰 숫자 강조
    "two_col": 3,   # slide4: 두 컬럼 (과제/기회)
    "vision":  4,   # slide5: 비전 (인용+원형 아이콘)
    "steps":   5,   # slide6: Step 카드 3개
    "metrics": 6,   # slide7: 성과 지표 (진행바)
    "chart":   7,   # slide8: 파이 차트
    "table":   8,   # slide9: 표
    "closing": 9,   # slide10: 결론
}


# ── 슬라이드 복제 ─────────────────────────────────────────────────────────────
def duplicate_slide(prs: Presentation, template_index: int):
    """템플릿 슬라이드를 복제해 프레젠테이션 끝에 추가. 복제된 슬라이드 반환."""
    src = prs.slides[template_index]

    # 새 슬라이드 추가 (같은 레이아웃 사용)
    slide_layout = src.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # spTree 전체를 src에서 복사
    src_tree = src.shapes._spTree
    new_tree = new_slide.shapes._spTree

    # 새 슬라이드의 기본 shape 제거 (nvGrpSpPr, grpSpPr 이후 모든 것)
    for child in list(new_tree)[2:]:
        new_tree.remove(child)

    # src 슬라이드의 모든 shape 복사 (nvGrpSpPr, grpSpPr 제외)
    for child in list(src_tree)[2:]:
        new_tree.append(copy.deepcopy(child))

    # 슬라이드 배경 복사 (bg 요소 있으면)
    src_cSld = src._element.find(qn_p("cSld"))
    new_cSld = new_slide._element.find(qn_p("cSld"))
    src_bg = src_cSld.find(qn_p("bg"))
    if src_bg is not None:
        existing_bg = new_cSld.find(qn_p("bg"))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        new_cSld.insert(0, copy.deepcopy(src_bg))

    # ★ 핵심: 소스 슬라이드의 관계(rels)를 새 슬라이드에 복사
    #   (이미지, 차트 등 rId2+ 참조가 유효하도록)
    SLIDE_LAYOUT_REL = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
    )
    for rId, rel in src.part.rels.items():
        if rel.reltype == SLIDE_LAYOUT_REL:
            continue  # slideLayout은 add_slide()가 이미 처리
        try:
            if rel.is_external:
                new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            else:
                new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception as e:
            pass  # 일부 rel 복사 실패 시 무시 (선택적 데코레이션)

    return new_slide


# ── 텍스트 조작 헬퍼 ──────────────────────────────────────────────────────────
def get_shape(slide, name: str):
    """이름으로 shape 찾기"""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def clear_tf(tf):
    """text_frame의 모든 단락/런 제거, 첫 단락만 남김"""
    txBody = tf._txBody
    paras = txBody.findall(qn_a("p"))
    for p in paras[1:]:
        txBody.remove(p)
    first_p = paras[0]
    for r in first_p.findall(qn_a("r")):
        first_p.remove(r)
    for br in first_p.findall(qn_a("br")):
        first_p.remove(br)
    return first_p


def get_ref_rpr(tf):
    """텍스트프레임에서 참조 런 속성(rPr) 추출"""
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn_a("rPr"))
            if rPr is not None:
                return copy.deepcopy(rPr)
    return None


def make_run(text: str, ref_rPr=None) -> etree._Element:
    """텍스트 런 XML 요소 생성"""
    r = etree.Element(qn_a("r"))
    if ref_rPr is not None:
        r.append(copy.deepcopy(ref_rPr))
    t = etree.SubElement(r, qn_a("t"))
    t.text = text
    return r


def make_para_with_text(text: str, ref_para=None, ref_rPr=None) -> etree._Element:
    """텍스트 포함 단락 XML 요소 생성"""
    if ref_para is not None:
        p = copy.deepcopy(ref_para)
        # 기존 런 제거
        for r in p.findall(qn_a("r")):
            p.remove(r)
        for br in p.findall(qn_a("br")):
            p.remove(br)
    else:
        p = etree.Element(qn_a("p"))

    if text:
        p.append(make_run(text, ref_rPr))
    return p


def set_shape_lines(slide, shape_name: str, lines):
    """shape를 이름으로 찾아 여러 줄 텍스트 설정 (서식 보존)"""
    shape = get_shape(slide, shape_name)
    if shape is None or not shape.has_text_frame:
        return False

    tf = shape.text_frame
    ref_rPr = get_ref_rpr(tf)
    txBody = tf._txBody
    first_p = txBody.findall(qn_a("p"))[0]
    ref_para = copy.deepcopy(first_p)

    # 기존 내용 지우기
    clear_tf(tf)
    first_p = txBody.findall(qn_a("p"))[0]

    if isinstance(lines, str):
        lines = [lines]

    # 첫 줄: 기존 첫 단락 재사용
    if lines:
        first_p.append(make_run(lines[0], ref_rPr))
    # 나머지 줄: 새 단락 추가
    for line in lines[1:]:
        p = make_para_with_text(line, ref_para=ref_para, ref_rPr=ref_rPr)
        txBody.append(p)
    return True


def set_shape_text(slide, shape_name: str, text: str):
    """단일 텍스트 설정"""
    return set_shape_lines(slide, shape_name, [text])


# ── 슬라이드 타입별 콘텐츠 주입 ────────────────────────────────────────────────
def fill_cover(slide, data: dict, config: dict):
    """표지 슬라이드 채우기"""
    # 빈 문자열은 config로 fallback
    year  = str(data.get("year")  or config.get("year",  "2025"))
    title = str(data.get("title") or config.get("title", "비즈니스 성장 전략"))
    sub   = str(data.get("subtitle") or config.get("subtitle", ""))

    set_shape_text(slide, "TextBox 6", year)
    set_shape_text(slide, "TextBox 3", title)
    set_shape_text(slide, "TextBox 1", sub)
    set_shape_text(slide, "TextBox 4", config.get("company", "Business proposal"))
    set_shape_text(slide, "TextBox 5", config.get("date", ""))


def fill_toc(slide, data: dict, config: dict):
    """목차 슬라이드 채우기"""
    items = data.get("items", [])
    nums  = [str(it.get("num", f"{i+1:02d}")) for i, it in enumerate(items)]
    texts = [str(it.get("text", "")) for it in items]

    set_shape_lines(slide, "TextBox 1", nums)
    set_shape_lines(slide, "TextBox 2", texts)
    set_shape_text(slide, "TextBox 4", config.get("company", "Business proposal"))
    set_shape_text(slide, "TextBox 5", config.get("date", ""))


def fill_common_header(slide, sec_num: str, title: str, config: dict):
    """공통 헤더 (섹션번호, 제목, 날짜, 회사명) 설정"""
    set_shape_text(slide, "TextBox 6", sec_num)
    set_shape_text(slide, "TextBox 3", title)
    set_shape_text(slide, "TextBox 4", config.get("company", "Business proposal"))
    set_shape_text(slide, "TextBox 5", config.get("date", ""))


def fill_stat(slide, data: dict, config: dict):
    """큰 숫자 강조 슬라이드 채우기 (slide3 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "01")),
        str(data.get("title", "")),
        config)

    set_shape_text(slide, "TextBox 7", str(data.get("big_number", "")))
    set_shape_text(slide, "TextBox 1", str(data.get("number_label", "")))

    bullets = []
    for i in range(1, 6):
        b = data.get(f"bullet{i}", "")
        if b:
            bullets.append(str(b))
    if bullets:
        set_shape_lines(slide, "TextBox 9", bullets)


def fill_two_col(slide, data: dict, config: dict):
    """두 컬럼 텍스트 슬라이드 채우기 (slide4 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "02")),
        str(data.get("title", "")),
        config)

    set_shape_text(slide, "TextBox 1", str(data.get("left_header", "")))
    left_items = [str(data.get(f"left{i}", "")) for i in range(1, 6) if data.get(f"left{i}")]
    set_shape_lines(slide, "TextBox 9", left_items)

    set_shape_text(slide, "TextBox 12", str(data.get("right_header", "")))
    right_items = [str(data.get(f"right{i}", "")) for i in range(1, 6) if data.get(f"right{i}")]
    set_shape_lines(slide, "TextBox 14", right_items)


def fill_vision(slide, data: dict, config: dict):
    """비전/인용 슬라이드 채우기 (slide5 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "03")),
        str(data.get("title", "")),
        config)

    quote = str(data.get("quote", data.get("body", "")))
    set_shape_lines(slide, "TextBox 1", quote.split("\n") if "\n" in quote else [quote])

    icons = []
    for i in range(1, 5):
        ic = data.get(f"icon{i}", "")
        if ic:
            icons.append(str(ic))
    for i, icon_text in enumerate(icons[:4]):
        # TextBox 9, 39, 42, 45 순서 (원래 템플릿 기준)
        names = ["TextBox 9", "TextBox 39", "TextBox 42", "TextBox 45"]
        if i < len(names):
            lines = icon_text.split("\n") if "\n" in icon_text else [icon_text]
            set_shape_lines(slide, names[i], lines)


def fill_steps(slide, data: dict, config: dict):
    """Step 카드 슬라이드 채우기 (slide6 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "04")),
        str(data.get("title", "")),
        config)

    # 각 스텝: step1_label, step1_title, step1_items (줄바꿈 구분)
    # 원본 shape: TextBox 12(label), TextBox 1(title), TextBox 9(items) — step1
    #             TextBox 29(label), TextBox 27(title), TextBox 28(items) — step2
    #             TextBox 34(label), TextBox 32(title), TextBox 33(items) — step3
    step_map = [
        ("TextBox 12", "TextBox 1",  "TextBox 9"),
        ("TextBox 29", "TextBox 27", "TextBox 28"),
        ("TextBox 34", "TextBox 32", "TextBox 33"),
    ]
    for i, (lbl_name, title_name, items_name) in enumerate(step_map, 1):
        lbl   = str(data.get(f"step{i}_label", f"step {i:02d}"))
        title = str(data.get(f"step{i}_title", ""))
        items_raw = str(data.get(f"step{i}_items", ""))
        items = [s.strip() for s in items_raw.split("\n") if s.strip()] if items_raw else []

        set_shape_text(slide, lbl_name, lbl)
        set_shape_text(slide, title_name, title)
        if items:
            set_shape_lines(slide, items_name, items)


def fill_metrics(slide, data: dict, config: dict):
    """성과 지표 슬라이드 채우기 (slide7 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "05")),
        str(data.get("title", "")),
        config)

    # 최대 4개 지표: metric1_label, metric1_value
    # TextBox 1/9(1번), 21/22(2번), 37/38(3번), 42/43(4번)
    label_names = ["TextBox 1", "TextBox 21", "TextBox 37", "TextBox 42"]
    value_names = ["TextBox 9", "TextBox 22", "TextBox 38", "TextBox 43"]

    for i, (ln, vn) in enumerate(zip(label_names, value_names), 1):
        lbl = str(data.get(f"metric{i}_label", ""))
        val = str(data.get(f"metric{i}_value", ""))
        if lbl:
            set_shape_text(slide, ln, lbl)
        if val:
            set_shape_text(slide, vn, val)


def fill_chart(slide, data: dict, config: dict):
    """파이 차트 슬라이드 채우기 (slide8 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "06")),
        str(data.get("title", "")),
        config)

    total_label = str(data.get("total_label", ""))
    set_shape_text(slide, "TextBox 1", total_label)

    # 범례 텍스트 업데이트
    legend_lines = []
    for i in range(1, 6):
        cat = data.get(f"cat{i}", "")
        val = data.get(f"val{i}", "")
        pct = data.get(f"pct{i}", "")
        if cat:
            if pct:
                legend_lines.append(f"{cat}: {val} ({pct}%)")
            else:
                legend_lines.append(f"{cat}: {val}")
    if legend_lines:
        set_shape_lines(slide, "TextBox 9", legend_lines)

    # 차트 데이터 업데이트 (내장 Excel 수정)
    _update_chart_data(slide, data)


def _update_chart_data(slide, data: dict):
    """슬라이드 내 차트의 내장 Excel 데이터 업데이트"""
    try:
        from pptx.oxml.ns import qn as pqn
        import openpyxl
        import io

        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            # 차트 데이터를 직접 XML로 업데이트
            plot = chart.plots[0]
            series = plot.series

            categories = []
            values = []
            for i in range(1, 6):
                cat = data.get(f"cat{i}", "")
                val = data.get(f"val{i}", "")
                pct = data.get(f"pct{i}", "")
                if cat and val:
                    categories.append(str(cat))
                    # 퍼센트 값 우선, 없으면 val 사용
                    try:
                        v = float(str(pct).replace("%", "")) / 100 if pct else float(str(val).replace("%", "").replace(",", ""))
                    except:
                        v = 0.25
                    values.append(v)

            if not categories:
                return

            # python-pptx ChartData로 교체
            from pptx.chart.data import ChartData
            chart_data = ChartData()
            chart_data.categories = categories
            chart_data.add_series("", values)
            chart.replace_data(chart_data)
            break
    except Exception as e:
        print(f"  [warn] 차트 데이터 업데이트 실패: {e}")


def fill_table(slide, data: dict, config: dict):
    """표 슬라이드 채우기 (slide9 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "07")),
        str(data.get("title", "")),
        config)

    # 표 shape 찾기
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows_data = data.get("rows", [])
            headers = data.get("headers", [])

            # 헤더 행 설정
            if headers and len(table.rows) > 0:
                for j, hdr in enumerate(headers):
                    if j < len(table.columns):
                        cell = table.cell(0, j)
                        if cell.text_frame.paragraphs:
                            _set_cell_text(cell, str(hdr))

            # 데이터 행 설정
            for i, row_data in enumerate(rows_data):
                tbl_row_idx = i + 1  # 0번 행 = 헤더
                if tbl_row_idx >= len(table.rows):
                    break
                for j, val in enumerate(row_data):
                    if j < len(table.columns):
                        cell = table.cell(tbl_row_idx, j)
                        _set_cell_text(cell, str(val))
            break


def _set_cell_text(cell, text: str):
    """표 셀 텍스트 설정"""
    tf = cell.text_frame
    ref_rPr = get_ref_rpr(tf)
    clear_tf(tf)
    txBody = tf._txBody
    first_p = txBody.findall(qn_a("p"))[0]
    first_p.append(make_run(text, ref_rPr))


def fill_closing(slide, data: dict, config: dict):
    """결론 슬라이드 채우기 (slide10 패턴)"""
    fill_common_header(slide,
        str(data.get("section_num", "08")),
        str(data.get("title", "결론")),
        config)

    body = str(data.get("body", data.get("body_text", "")))
    if body:
        lines = body.split("\n") if "\n" in body else [body]
        set_shape_lines(slide, "TextBox 7", lines)

    dl_text = str(data.get("download_text", "사용폰트 다운로드 링크"))
    if dl_text:
        set_shape_text(slide, "TextBox 10", dl_text)


# ── 필러 디스패치 테이블 ───────────────────────────────────────────────────────
FILLERS = {
    "cover":   fill_cover,
    "toc":     fill_toc,
    "stat":    fill_stat,
    "two_col": fill_two_col,
    "vision":  fill_vision,
    "steps":   fill_steps,
    "metrics": fill_metrics,
    "chart":   fill_chart,
    "table":   fill_table,
    "closing": fill_closing,
}


# ── ZIP 레벨 슬라이드 정리 ────────────────────────────────────────────────────
def _remove_slides(prs: Presentation, indices: list):
    """
    sldIdLst에서만 참조 제거 (저장 전 마킹).
    실제 ZIP 정리는 _clean_pptx_zip()으로 저장 후 수행.
    """
    sldIdLst = prs.slides._sldIdLst
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for idx in sorted(indices, reverse=True):
        all_sld = list(sldIdLst)
        if idx >= len(all_sld):
            continue
        sId = all_sld[idx]
        sldIdLst.remove(sId)


def _clean_pptx_zip(pptx_path: str):
    """
    저장된 PPTX ZIP에서 sldIdLst에 없는 슬라이드 파일 및 관계를 완전 제거.
    PowerPoint 호환성을 위한 필수 정리 작업.
    """
    import zipfile

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
    NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"

    with zipfile.ZipFile(pptx_path, "r") as z:
        names = z.namelist()
        contents = {n: z.read(n) for n in names}

    # presentation.xml 파싱 → 실제 참조되는 slide rId 목록
    prs_xml = etree.fromstring(contents["ppt/presentation.xml"])
    sldIdLst = prs_xml.find(f"{{{NS_P}}}sldIdLst")
    active_rids = set()
    if sldIdLst is not None:
        for sId in sldIdLst:
            rId = sId.get(f"{{{NS_R}}}id")
            if rId:
                active_rids.add(rId)

    # presentation.xml.rels 파싱 → 참조 안 된 slide rel 찾기
    rels_key = "ppt/_rels/presentation.xml.rels"
    rels_xml = etree.fromstring(contents[rels_key])
    slide_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    orphan_targets = []

    for rel in list(rels_xml):
        if rel.get("Type") == slide_type and rel.get("Id") not in active_rids:
            orphan_targets.append(rel.get("Target", ""))
            rels_xml.remove(rel)

    if not orphan_targets:
        return  # 정리할 것 없음

    contents[rels_key] = etree.tostring(
        rels_xml, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    # 고아 슬라이드 파일 목록
    orphan_files = set()
    for target in orphan_targets:
        slide_path = f"ppt/{target}"
        orphan_files.add(slide_path)
        slide_name = target.split("/")[-1]
        orphan_files.add(f"ppt/slides/_rels/{slide_name}.rels")

    # [Content_Types].xml에서 고아 슬라이드 제거
    ct_xml = etree.fromstring(contents["[Content_Types].xml"])
    for override in list(ct_xml):
        partname = override.get("PartName", "")
        for f in orphan_files:
            if partname == f"/{f}" or partname == f.replace("ppt/", "/ppt/", 1):
                ct_xml.remove(override)
                break
    contents["[Content_Types].xml"] = etree.tostring(
        ct_xml, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    # 정리된 ZIP 다시 쓰기
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in contents.items():
            if name not in orphan_files:
                z.writestr(name, data)

    print(f"[build_pptx] ZIP 정리 완료: {len(orphan_targets)}개 고아 슬라이드 제거")


# ── 메인 빌드 함수 ────────────────────────────────────────────────────────────
def build_presentation(slide_plan: list, config: dict, template_path: str, output_path: str):
    """
    slide_plan: 슬라이드 목록 [{"type": ..., "data": {...}}, ...]
    config: 공통 설정 (company, date, ...)
    """
    print(f"[build_pptx] Loading template: {template_path}")
    prs = Presentation(template_path)

    # 원본 템플릿 슬라이드 수 기록
    n_template = len(prs.slides)
    print(f"[build_pptx] Template has {n_template} slides")

    # 각 슬라이드 계획대로 복제 + 콘텐츠 주입
    new_slide_indices = []
    for slide_spec in slide_plan:
        stype = slide_spec.get("type", "content")
        tmpl_idx = TEMPLATE.get(stype, TEMPLATE.get("two_col", 3))

        print(f"  → Slide {slide_spec.get('order', '?')}: type={stype} template_idx={tmpl_idx}")

        new_slide = duplicate_slide(prs, tmpl_idx)
        new_slide_indices.append(len(prs.slides) - 1)

        # 콘텐츠 주입
        filler = FILLERS.get(stype)
        if filler:
            slide_data = slide_spec.get("data", slide_spec)
            # section_num, title을 data에 병합
            if "section_num" in slide_spec:
                slide_data = dict(slide_data)
                slide_data.setdefault("section_num", slide_spec["section_num"])
            if "title" in slide_spec:
                slide_data = dict(slide_data)
                slide_data.setdefault("title", slide_spec["title"])
            if "items" in slide_spec:
                slide_data = dict(slide_data)
                slide_data["items"] = slide_spec["items"]
            try:
                filler(new_slide, slide_data, config)
            except Exception as e:
                print(f"    [warn] 슬라이드 채우기 실패: {e}")

    # 원본 템플릿 슬라이드 제거 (앞 n_template 개)
    print(f"[build_pptx] Removing {n_template} template slides...")
    _remove_slides(prs, list(range(n_template)))

    print(f"[build_pptx] Final slide count: {len(prs.slides)}")
    prs.save(output_path)
    print(f"[build_pptx] Saved: {output_path}")

    # ZIP 레벨 정리 (PowerPoint 호환성)
    _clean_pptx_zip(output_path)

    return output_path


# ── CLI 진입점 ────────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="스톰데이즈 템플릿 기반 PPTX 생성")
    parser.add_argument("slide_plan_json", help="slide_plan.json 경로")
    parser.add_argument("output_pptx", help="출력 PPTX 경로")
    parser.add_argument("--template", default=None, help="템플릿 PPTX 경로")
    args = parser.parse_args()

    # 기본 템플릿 경로: 스크립트 기준 ../template/stormdays.pptx
    template_path = args.template or str(
        Path(__file__).parent.parent / "template" / "stormdays.pptx"
    )

    with open(args.slide_plan_json, encoding="utf-8") as f:
        plan_data = json.load(f)

    config = plan_data.get("config", {})
    slide_plan = plan_data.get("slide_plan", [])

    if not slide_plan:
        print("[build_pptx] ERROR: slide_plan이 비어있습니다.")
        sys.exit(1)

    build_presentation(slide_plan, config, template_path, args.output_pptx)


if __name__ == "__main__":
    main()
